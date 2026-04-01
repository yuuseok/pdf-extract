import os
import tempfile
from pptx import Presentation
from pptx.util import Inches
from app.service.pptx_service import PptxService


def _create_test_pptx(path: str):
    prs = Presentation()
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "발표 제목"
    slide.placeholders[1].text = "발표 내용입니다."
    # Slide 2
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "두 번째 슬라이드"
    slide2.placeholders[1].text = "추가 내용입니다."
    prs.save(path)


def test_pptx_extract_returns_three_formats():
    svc = PptxService()
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        _create_test_pptx(f.name)
        result = svc.extract(f.name)
    os.unlink(f.name)

    assert "text" in result
    assert "markdown" in result
    assert "json" in result
    assert "발표 제목" in result["text"]
    assert result["page_count"] == 2


def test_pptx_slide_structure():
    svc = PptxService()
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        _create_test_pptx(f.name)
        result = svc.extract(f.name)
    os.unlink(f.name)

    slides = result["json"]
    assert len(slides) == 2
    assert slides[0]["slide_number"] == 1
    assert slides[1]["slide_number"] == 2


def test_pptx_markdown_slides():
    svc = PptxService()
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as f:
        _create_test_pptx(f.name)
        result = svc.extract(f.name)
    os.unlink(f.name)

    assert "## Slide 1:" in result["markdown"]
    assert "## Slide 2:" in result["markdown"]
