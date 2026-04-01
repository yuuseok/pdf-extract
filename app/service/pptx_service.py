import logging
from pptx import Presentation

logger = logging.getLogger(__name__)


class PptxService:
    """PowerPoint(.pptx) 파일에서 text, markdown, json을 추출."""

    def extract(self, file_path: str) -> dict:
        prs = Presentation(file_path)
        elements = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                "type": "slide",
                "slide_number": slide_idx,
                "title": None,
                "contents": [],
            }

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if not text:
                        continue

                    # 제목 shape 감지
                    if shape.shape_type and hasattr(shape, "placeholder_format"):
                        if shape.placeholder_format and shape.placeholder_format.idx == 0:
                            slide_data["title"] = text
                            continue

                    slide_data["contents"].append({
                        "type": "text",
                        "content": text,
                    })

                if shape.has_table:
                    table_data = self._extract_table(shape.table)
                    if table_data["rows"]:
                        slide_data["contents"].append(table_data)

            # 제목이 없으면 첫 텍스트를 제목으로
            if not slide_data["title"] and slide_data["contents"]:
                for c in slide_data["contents"]:
                    if c["type"] == "text":
                        slide_data["title"] = c["content"]
                        slide_data["contents"].remove(c)
                        break

            elements.append(slide_data)

        text = self._to_text(elements)
        markdown = self._to_markdown(elements)

        return {
            "text": text,
            "markdown": markdown,
            "json": elements,
            "json_raw": elements,
            "page_count": len(elements),
        }

    def _extract_table(self, table) -> dict:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        return {"type": "table", "rows": rows}

    def _to_text(self, elements: list) -> str:
        lines = []
        for slide in elements:
            if slide.get("title"):
                lines.append(slide["title"])
            for content in slide.get("contents", []):
                if content["type"] == "text":
                    lines.append(content["content"])
                elif content["type"] == "table":
                    for row in content["rows"]:
                        lines.append("\t".join(row))
            lines.append("")
        return "\n".join(lines)

    def _to_markdown(self, elements: list) -> str:
        lines = []
        for slide in elements:
            slide_num = slide.get("slide_number", "")
            title = slide.get("title", f"Slide {slide_num}")
            lines.append(f"## Slide {slide_num}: {title}")
            lines.append("")

            for content in slide.get("contents", []):
                if content["type"] == "text":
                    lines.append(content["content"])
                    lines.append("")
                elif content["type"] == "table":
                    rows = content["rows"]
                    if not rows:
                        continue
                    col_count = max(len(r) for r in rows)
                    header = rows[0] + [""] * (col_count - len(rows[0]))
                    lines.append("| " + " | ".join(header) + " |")
                    lines.append("| " + " | ".join(["---"] * col_count) + " |")
                    for row in rows[1:]:
                        padded = row + [""] * (col_count - len(row))
                        lines.append("| " + " | ".join(padded) + " |")
                    lines.append("")
        return "\n".join(lines)
